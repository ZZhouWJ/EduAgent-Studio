import tempfile
import unittest
from pathlib import Path
from unittest.mock import Mock, patch

from docx import Document
from pptx import Presentation

from app.rag.parser import parse_document_file


class DocumentParserTests(unittest.TestCase):
    def test_parses_utf8_text_file(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "lesson.txt"
            path.write_text("事务的原子性保证操作全部成功或全部回滚。", encoding="utf-8")

            chunks = parse_document_file(str(path), "text")

        self.assertEqual(len(chunks), 1)
        self.assertIn("原子性", chunks[0]["content"])

    def test_parses_docx_headings_paragraphs_and_tables(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "lesson.docx"
            document = Document()
            document.add_heading("数据库事务", level=2)
            document.add_paragraph("事务包含四项 ACID 特性。")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "原子性"
            table.cell(0, 1).text = "全部成功或全部回滚"
            document.save(path)

            chunks = parse_document_file(str(path), "word")

        content = "\n".join(chunk["content"] for chunk in chunks)
        self.assertIn("ACID", content)
        self.assertIn("全部回滚", content)

    def test_parses_pptx_by_slide_and_preserves_page_number(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "lesson.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "隔离级别"
            slide.placeholders[1].text = "读未提交\n读已提交\n可重复读\n串行化"
            presentation.save(path)

            chunks = parse_document_file(str(path), "ppt")

        self.assertEqual(chunks[0]["source_page"], 1)
        self.assertEqual(chunks[0]["title"], "隔离级别")
        self.assertIn("串行化", chunks[0]["content"])

    @patch("pypdf.PdfReader")
    def test_parses_pdf_by_page(self, reader: Mock):
        first_page = Mock()
        first_page.extract_text.return_value = "第一页面介绍事务。"
        second_page = Mock()
        second_page.extract_text.return_value = "第二页面介绍锁。"
        reader.return_value.pages = [first_page, second_page]

        chunks = parse_document_file("lesson.pdf", "pdf")

        self.assertEqual([chunk["source_page"] for chunk in chunks], [1, 2])


if __name__ == "__main__":
    unittest.main()
