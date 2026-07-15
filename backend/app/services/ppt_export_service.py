"""将结构化课件内容导出为可继续编辑的 PowerPoint 文件。"""

from io import BytesIO
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def build_presentation(title: str, slides: List[Dict[str, Any]]) -> BytesIO:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = "EduAgent Studio 个性化学习课件"
    _style_title(cover.shapes.title.text_frame, 30)
    _style_subtitle(cover.placeholders[1].text_frame)

    for index, item in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = item.get("title") or f"第 {index} 页"
        _style_title(slide.shapes.title.text_frame, 26)

        body = slide.placeholders[1].text_frame
        body.clear()
        points = item.get("points") or item.get("bullets") or []
        for point_index, point in enumerate(points):
            paragraph = body.paragraphs[0] if point_index == 0 else body.add_paragraph()
            paragraph.text = str(point)
            paragraph.level = 0
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(30, 41, 59)
            paragraph.space_after = Pt(10)

        notes = str(item.get("notes") or "").strip()
        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = notes

        page_box = slide.shapes.add_textbox(
            presentation.slide_width - Inches(1.1),
            presentation.slide_height - Inches(0.45),
            Inches(0.7),
            Inches(0.25),
        )
        page = page_box.text_frame.paragraphs[0]
        page.text = str(index)
        page.alignment = PP_ALIGN.RIGHT
        page.font.name = "Microsoft YaHei"
        page.font.size = Pt(10)
        page.font.color.rgb = RGBColor(100, 116, 139)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output


def _style_title(text_frame, size: int) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(15, 23, 42)


def _style_subtitle(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(71, 85, 105)
