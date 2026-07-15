"""
文档解析器模块。

支持解析 PDF、DOCX、PPTX、Markdown 和纯文本文件，生成结构化 chunks。
提供 BM25 关键词提取功能。
"""

import re
from typing import Any, Dict, List


def extract_bm25_terms(text: str) -> List[str]:
    """
    从文本中提取 BM25 检索关键词。

    使用简单中文分词 + 停用词过滤。

    Args:
        text: 输入文本

    Returns:
        关键词列表（去重、按长度排序）
    """
    STOP_WORDS = {
        "的", "了", "在", "是", "我", "有", "和", "就",
        "不", "人", "都", "一", "一个", "上", "也", "很",
        "到", "说", "要", "去", "你", "会", "着", "没有",
        "看", "好", "自己", "这", "那", "它", "什么",
        "可以", "这个", "那个", "因为", "所以", "但是",
        "如果", "或者", "而且", "还是", "只是", "还有",
        "然后", "这样", "那样", "如何", "怎么", "为什么",
        "哪", "哪些", "什么", "如何", "是否", "是不是",
    }

    # 提取中文词和英文词
    words = re.findall(r"[一-鿿]+|[a-zA-Z0-9]+", text.lower())
    # 过滤停用词、长度小于2的词
    filtered = [w for w in words if w not in STOP_WORDS and len(w) >= 2]
    # 去重并保持顺序
    seen = set()
    unique = []
    for w in filtered:
        if w not in seen:
            seen.add(w)
            unique.append(w)
    return unique


def _tokenize_for_chunking(text: str) -> List[str]:
    """
    文档切分用的分词器，保留更多短词。
    """
    STOP_WORDS = {
        "的", "了", "在", "是", "我", "有", "和", "就",
        "不", "人", "都", "一", "一个", "上", "也", "很",
    }
    words = re.findall(r"[一-鿿]+|[a-zA-Z0-9]+", text.lower())
    return [w for w in words if w not in STOP_WORDS and len(w) >= 2]


def parse_markdown(content: str) -> List[Dict[str, Any]]:
    """
    解析 Markdown 文档，生成 chunks。

    按标题（##）切分章节，每章节为一个 chunk。
    如果章节过长（超过2000字符），进一步按段落切分。

    Args:
        content: Markdown 文本内容

    Returns:
        List[Dict]: chunks 列表，每项包含 title, content, source_paragraph
    """
    chunks = []

    # 按 ## 标题分割
    sections = re.split(r"\n(?=##\s)", content)

    for idx, section in enumerate(sections):
        section = section.strip()
        if not section:
            continue

        # 提取标题（第一行的 ## 标题）
        title_match = re.match(r"^(#{1,6})\s+(.+)$", section.split("\n")[0])
        if title_match:
            title = title_match.group(2).strip()
            body = "\n".join(section.split("\n")[1:]).strip()
        else:
            # 无标题，使用"第N段"作为标题
            title = f"第 {idx + 1} 节"
            body = section

        if not body:
            continue

        # 如果内容过长，按段落切分
        if len(body) > 2000:
            paragraphs = re.split(r"\n\n+", body)
            current_chunk = ""
            current_para_idx = 0

            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue

                if len(current_chunk) + len(para) + 2 <= 2000:
                    current_chunk += para + "\n\n"
                else:
                    if current_chunk:
                        chunks.append({
                            "title": title if current_para_idx == 0 else f"{title} (续)",
                            "content": current_chunk.strip(),
                            "source_paragraph": current_para_idx,
                        })
                    current_chunk = para + "\n\n"
                    current_para_idx += 1

            if current_chunk:
                chunks.append({
                    "title": title if current_para_idx == 0 else f"{title} (续)",
                    "content": current_chunk.strip(),
                    "source_paragraph": current_para_idx,
                })
        else:
            chunks.append({
                "title": title,
                "content": body,
                "source_paragraph": 0,
            })

    return chunks


def parse_text(content: str) -> List[Dict[str, Any]]:
    """
    解析纯文本文件，生成 chunks。

    按段落切分，每段为一个 chunk。
    过短的段落（少于100字符）会与下一段合并。

    Args:
        content: 纯文本内容

    Returns:
        List[Dict]: chunks 列表
    """
    chunks = []

    # 按一个或多个空行分割段落
    paragraphs = re.split(r"\n\s*\n", content)

    current_chunk = ""
    chunk_index = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # 如果当前 chunk 为空，直接放入
        if not current_chunk:
            current_chunk = para
            continue

        # 如果合并后长度合适，合并
        if len(current_chunk) + len(para) + 2 <= 2000:
            current_chunk += "\n\n" + para
        else:
            # 当前 chunk 足够大，保存并开始新的
            chunks.append({
                "title": f"段落 {chunk_index + 1}",
                "content": current_chunk.strip(),
                "source_paragraph": chunk_index,
            })
            chunk_index += 1
            current_chunk = para

    # 处理最后一个 chunk
    if current_chunk:
        chunks.append({
            "title": f"段落 {chunk_index + 1}",
            "content": current_chunk.strip(),
            "source_paragraph": chunk_index,
        })

    return chunks


def parse_document(content: str, file_type: str) -> List[Dict[str, Any]]:
    """
    通用文档解析入口，根据文件类型选择解析方法。

    Args:
        content: 文档内容
        file_type: 文件类型 (markdown/md/txt/text)

    Returns:
        List[Dict]: chunks 列表
    """
    file_type = file_type.lower()

    if file_type in ("markdown", "md"):
        return parse_markdown(content)
    elif file_type in ("text", "txt"):
        return parse_text(content)
    else:
        return parse_text(content)


def parse_document_file(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    """从真实文件格式提取文本，并生成带页码来源的知识片段。"""
    normalized_type = file_type.lower()
    if normalized_type in {"markdown", "md", "text", "txt"}:
        with open(file_path, "r", encoding="utf-8") as file:
            return parse_document(file.read(), normalized_type)
    if normalized_type == "pdf":
        return _parse_pdf_file(file_path)
    if normalized_type == "word":
        return _parse_docx_file(file_path)
    if normalized_type == "ppt":
        return _parse_pptx_file(file_path)
    raise ValueError(f"不支持的文档类型: {normalized_type}")


def _parse_pdf_file(file_path: str) -> List[Dict[str, Any]]:
    from pypdf import PdfReader

    chunks: List[Dict[str, Any]] = []
    reader = PdfReader(file_path)
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        page_chunks = parse_text(text)
        for chunk in page_chunks:
            chunk["title"] = f"第 {page_number} 页 - {chunk['title']}"
            chunk["source_page"] = page_number
            chunks.append(chunk)
    return chunks


def _parse_docx_file(file_path: str) -> List[Dict[str, Any]]:
    from docx import Document

    document = Document(file_path)
    lines: List[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name.lower() if paragraph.style else ""
        if style_name.startswith("heading"):
            level_match = re.search(r"(\d+)", style_name)
            level = min(6, int(level_match.group(1))) if level_match else 2
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))

    content = "\n\n".join(lines)
    if not content:
        return []
    return parse_markdown(content) if any(line.startswith("#") for line in lines) else parse_text(content)


def _parse_pptx_file(file_path: str) -> List[Dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    chunks: List[Dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines: List[str] = []
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))

        content = "\n\n".join(lines).strip()
        if not title and not content:
            continue
        chunks.append({
            "title": title or f"第 {slide_number} 页",
            "content": content or title,
            "source_page": slide_number,
            "source_paragraph": 0,
        })
    return chunks
