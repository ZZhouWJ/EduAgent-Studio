#!/usr/bin/env python3
"""Build the contest DOCX and PPTX deliverables from reviewed Markdown sources."""

from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import Iterable, Sequence

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
SUBMISSION = ROOT / "submission"
SCREENSHOTS = SUBMISSION / "assets" / "screenshots"

CHARCOAL = "17202A"
INK = "202B37"
TEAL = "0F766E"
TEAL_LIGHT = "DDF4F0"
BLUE = "2563EB"
BLUE_LIGHT = "EAF1FF"
CORAL = "DB6257"
CORAL_LIGHT = "FCE9E6"
MUTED = "607080"
LINE = "D9E1E7"
PAPER = "F5F8F9"
WHITE = "FFFFFF"

PPT_W = 13.333
PPT_H = 7.5
CJK_FONT = "PingFang SC"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def ppt_rgb(hex_value: str) -> PptRGBColor:
    return PptRGBColor.from_string(hex_value)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_margins(cell, top=90, start=100, bottom=90, end=100) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def set_repeat_table_header(row) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), "true")
    tr_pr.append(tbl_header)


def set_keep_with_next(paragraph, value: bool = True) -> None:
    paragraph.paragraph_format.keep_with_next = value


def set_font(run, name: str = CJK_FONT, size: float | None = None) -> None:
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size is not None:
        run.font.size = Pt(size)


def add_page_number(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = paragraph.add_run("第 ")
    set_font(run, size=8.5)
    fld_char_1 = OxmlElement("w:fldChar")
    fld_char_1.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = " PAGE "
    fld_char_2 = OxmlElement("w:fldChar")
    fld_char_2.set(qn("w:fldCharType"), "end")
    run._r.extend([fld_char_1, instr_text, fld_char_2])
    tail = paragraph.add_run(" 页")
    set_font(tail, size=8.5)
    run.font.color.rgb = rgb(MUTED)
    tail.font.color.rgb = rgb(MUTED)


def configure_doc_styles(document: Document) -> None:
    normal = document.styles["Normal"]
    normal.font.name = CJK_FONT
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), CJK_FONT)
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = rgb(INK)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.35

    heading_specs = {
        "Title": (28, CHARCOAL, 0, 14),
        "Subtitle": (13, MUTED, 0, 8),
        "Heading 1": (20, CHARCOAL, 18, 9),
        "Heading 2": (15, TEAL, 14, 7),
        "Heading 3": (12, INK, 10, 5),
        "Heading 4": (10.5, MUTED, 8, 4),
    }
    for style_name, (size, color, before, after) in heading_specs.items():
        style = document.styles[style_name]
        style.font.name = CJK_FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), CJK_FONT)
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = rgb(color)
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)
        style.paragraph_format.keep_with_next = True

    for style_name in ("List Bullet", "List Number"):
        style = document.styles[style_name]
        style.font.name = CJK_FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), CJK_FONT)
        style.font.size = Pt(10.2)
        style.paragraph_format.left_indent = Cm(0.65)
        style.paragraph_format.first_line_indent = Cm(-0.35)
        style.paragraph_format.space_after = Pt(3)


def add_inline_runs(paragraph, text: str) -> None:
    token_pattern = re.compile(r"(\*\*.+?\*\*|`[^`]+`|\[[^\]]+\]\([^)]+\))")
    position = 0
    for match in token_pattern.finditer(text):
        if match.start() > position:
            run = paragraph.add_run(text[position:match.start()])
            set_font(run)
        token = match.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            run.bold = True
            run.font.color.rgb = rgb(CHARCOAL)
            set_font(run)
        elif token.startswith("`"):
            run = paragraph.add_run(token[1:-1])
            run.font.name = "Menlo"
            run._element.rPr.rFonts.set(qn("w:eastAsia"), CJK_FONT)
            run.font.size = Pt(9)
            run.font.color.rgb = rgb(BLUE)
            shading = OxmlElement("w:shd")
            shading.set(qn("w:fill"), BLUE_LIGHT)
            run._r.get_or_add_rPr().append(shading)
        else:
            link_match = re.match(r"\[([^\]]+)\]\(([^)]+)\)", token)
            label = link_match.group(1) if link_match else token
            run = paragraph.add_run(label)
            run.font.color.rgb = rgb(BLUE)
            run.underline = True
            set_font(run)
        position = match.end()
    if position < len(text):
        run = paragraph.add_run(text[position:])
        set_font(run)


def add_doc_cover(document: Document, title: str, category: str) -> None:
    for _ in range(4):
        document.add_paragraph()
    kicker = document.add_paragraph()
    kicker.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = kicker.add_run("第十五届中国软件杯 · A3 赛题")
    run.bold = True
    run.font.color.rgb = rgb(TEAL)
    set_font(run, size=11)

    title_p = document.add_paragraph(style="Title")
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.add_run(title)
    subtitle = document.add_paragraph(style="Subtitle")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("智学工坊 EduAgent Studio")

    document.add_paragraph()
    table = document.add_table(rows=4, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    facts = [
        ("交付类型", category),
        ("赛题方向", "基于大模型的个性化资源生成与学习多智能体系统"),
        ("文档版本", "V2.1"),
        ("核对日期", "2026-07-17"),
    ]
    for row, (label, value) in zip(table.rows, facts):
        row.cells[0].width = Cm(3.2)
        row.cells[1].width = Cm(11.2)
        set_cell_shading(row.cells[0], TEAL_LIGHT)
        set_cell_shading(row.cells[1], PAPER)
        for cell in row.cells:
            set_cell_margins(cell, 130, 150, 130, 150)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        p1 = row.cells[0].paragraphs[0]
        p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1 = p1.add_run(label)
        r1.bold = True
        r1.font.color.rgb = rgb(TEAL)
        set_font(r1, size=9.5)
        p2 = row.cells[1].paragraphs[0]
        r2 = p2.add_run(value)
        set_font(r2, size=9.5)

    document.add_paragraph()
    note = document.add_paragraph()
    note.alignment = WD_ALIGN_PARAGRAPH.CENTER
    note_run = note.add_run("以当前 main 分支代码、数据库迁移、自动化测试和真实运行结果为准")
    note_run.italic = True
    note_run.font.color.rgb = rgb(MUTED)
    set_font(note_run, size=9)
    document.add_page_break()


def is_table_separator(line: str) -> bool:
    cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
    return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells)


def split_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def add_doc_table(document: Document, rows: Sequence[Sequence[str]]) -> None:
    if not rows:
        return
    column_count = max(len(row) for row in rows)
    table = document.add_table(rows=len(rows), cols=column_count)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    table.autofit = True
    for row_index, source_row in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            set_cell_margins(cell)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            value = source_row[column_index] if column_index < len(source_row) else ""
            paragraph = cell.paragraphs[0]
            paragraph.paragraph_format.space_after = Pt(0)
            add_inline_runs(paragraph, value)
            if row_index == 0:
                set_cell_shading(cell, CHARCOAL)
                for run in paragraph.runs:
                    run.font.color.rgb = rgb(WHITE)
                    run.bold = True
            elif row_index % 2 == 0:
                set_cell_shading(cell, PAPER)
        table.rows[row_index]._tr.get_or_add_trPr().append(OxmlElement("w:cantSplit"))
    set_repeat_table_header(table.rows[0])
    document.add_paragraph().paragraph_format.space_after = Pt(2)


def add_code_block(document: Document, lines: Sequence[str]) -> None:
    table = document.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0)
    set_cell_shading(cell, "EEF2F4")
    set_cell_margins(cell, 140, 180, 140, 180)
    paragraph = cell.paragraphs[0]
    paragraph.paragraph_format.space_after = Pt(0)
    paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
    run = paragraph.add_run("\n".join(lines))
    run.font.name = "Menlo"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), CJK_FONT)
    run.font.size = Pt(8.5)
    run.font.color.rgb = rgb(INK)
    document.add_paragraph().paragraph_format.space_after = Pt(2)


def add_quote_block(document: Document, text: str) -> None:
    table = document.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    set_cell_shading(cell, TEAL_LIGHT)
    set_cell_margins(cell, 100, 180, 100, 180)
    paragraph = cell.paragraphs[0]
    paragraph.paragraph_format.space_after = Pt(0)
    add_inline_runs(paragraph, text)
    for run in paragraph.runs:
        run.font.color.rgb = rgb(TEAL)
        run.italic = True


def add_markdown_body(document: Document, markdown_text: str) -> None:
    lines = markdown_text.splitlines()
    index = 0
    in_code = False
    code_lines: list[str] = []
    paragraph_lines: list[str] = []

    def flush_paragraph() -> None:
        if not paragraph_lines:
            return
        text = " ".join(part.strip() for part in paragraph_lines).strip()
        paragraph_lines.clear()
        if text:
            paragraph = document.add_paragraph()
            add_inline_runs(paragraph, text)

    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if stripped.startswith("```"):
            flush_paragraph()
            if in_code:
                add_code_block(document, code_lines)
                code_lines = []
                in_code = False
            else:
                in_code = True
            index += 1
            continue
        if in_code:
            code_lines.append(line)
            index += 1
            continue

        if stripped.startswith("|") and index + 1 < len(lines) and is_table_separator(lines[index + 1]):
            flush_paragraph()
            table_rows = [split_table_row(stripped)]
            index += 2
            while index < len(lines) and lines[index].strip().startswith("|"):
                table_rows.append(split_table_row(lines[index]))
                index += 1
            add_doc_table(document, table_rows)
            continue

        heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
        if heading:
            flush_paragraph()
            level = len(heading.group(1))
            if level == 1:
                index += 1
                continue
            paragraph = document.add_paragraph(style=f"Heading {level - 1}")
            add_inline_runs(paragraph, heading.group(2))
            index += 1
            continue

        if stripped.startswith(">"):
            flush_paragraph()
            quote_lines = []
            while index < len(lines) and (lines[index].strip().startswith(">") or not lines[index].strip()):
                if lines[index].strip().startswith(">"):
                    quote_lines.append(lines[index].strip()[1:].strip())
                index += 1
            if quote_lines:
                add_quote_block(document, " · ".join(quote_lines))
            continue

        bullet = re.match(r"^[-*]\s+(.+)$", stripped)
        numbered = re.match(r"^\d+\.\s+(.+)$", stripped)
        if bullet or numbered:
            flush_paragraph()
            paragraph = document.add_paragraph(style="List Bullet" if bullet else "List Number")
            add_inline_runs(paragraph, (bullet or numbered).group(1))
            index += 1
            continue

        if not stripped:
            flush_paragraph()
            index += 1
            continue

        paragraph_lines.append(stripped)
        index += 1

    flush_paragraph()
    if code_lines:
        add_code_block(document, code_lines)


def build_docx(source: Path, output: Path, category: str) -> None:
    markdown_text = source.read_text(encoding="utf-8")
    title_match = re.search(r"^#\s+(.+)$", markdown_text, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else source.stem

    document = Document()
    section = document.sections[0]
    section.top_margin = Cm(2.2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.2)
    section.right_margin = Cm(2)
    section.header_distance = Cm(1)
    section.footer_distance = Cm(1)
    configure_doc_styles(document)

    header = section.header.paragraphs[0]
    header.text = "EDUAGENT STUDIO  ·  A3 CONTEST DELIVERY"
    header_run = header.runs[0]
    header_run.font.color.rgb = rgb(MUTED)
    set_font(header_run, name="Aptos", size=8)
    add_page_number(section.footer.paragraphs[0])

    add_doc_cover(document, title, category)
    add_markdown_body(document, markdown_text)

    core = document.core_properties
    core.title = title
    core.subject = "EduAgent Studio 中国软件杯 A3 参赛材料"
    core.author = "EduAgent Studio 参赛团队"
    core.keywords = "EduAgent Studio, A3, 多智能体, 个性化学习"
    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)


def ppt_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = CJK_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.06,
):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = PptInches(margin)
    frame.margin_right = PptInches(margin)
    frame.margin_top = PptInches(margin)
    frame.margin_bottom = PptInches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = PptPt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    return box


def ppt_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(fill)
    shape.line.color.rgb = ppt_rgb(line)
    shape.line.width = PptPt(1)
    return shape


def ppt_line(slide, x1, y1, x2, y2, color=LINE, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        PptInches(x1),
        PptInches(y1),
        PptInches(x2),
        PptInches(y2),
    )
    line.line.color.rgb = ppt_rgb(color)
    line.line.width = PptPt(width)
    return line


def add_slide_chrome(slide, number: int, title: str, kicker: str) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = ppt_rgb(PAPER)
    ppt_rect(slide, 0, 0, PPT_W, 0.1, fill=TEAL, line=TEAL, radius=False)
    ppt_text(slide, kicker.upper(), 0.65, 0.36, 5.8, 0.28, 9, TEAL, True, font="Aptos")
    ppt_text(slide, title, 0.65, 0.68, 11.8, 0.55, 26, CHARCOAL, True)
    ppt_line(slide, 0.65, 1.34, 12.68, 1.34, LINE, 1)
    ppt_text(slide, "EduAgent Studio · A3", 0.65, 7.14, 3.8, 0.2, 8, MUTED, font="Aptos")
    ppt_text(slide, f"{number:02d}", 12.15, 7.08, 0.5, 0.25, 9, MUTED, True, font="Aptos", align=PP_ALIGN.RIGHT)


def add_bullets(slide, items: Sequence[str], x, y, w, h, size=16, color=INK, accent=TEAL):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = CJK_FONT
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = ppt_rgb(color)
        paragraph.space_after = PptPt(10)
        paragraph.level = 0
        paragraph.text = f"●  {item}"
        for run in paragraph.runs:
            run.font.name = CJK_FONT
            run.font.size = PptPt(size)
            run.font.color.rgb = ppt_rgb(color)
    return box


def add_picture_crop(slide, path: Path, x, y, w, h):
    if not path.exists():
        placeholder = ppt_rect(slide, x, y, w, h, fill=WHITE, line=CORAL)
        ppt_text(slide, f"缺少截图\n{path.name}", x + 0.2, y + h / 2 - 0.35, w - 0.4, 0.7, 13, CORAL, True, align=PP_ALIGN.CENTER)
        return placeholder
    with Image.open(path) as image:
        image_w, image_h = image.size
    picture = slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    image_ratio = image_w / image_h
    target_ratio = w / h
    if image_ratio > target_ratio:
        visible = target_ratio / image_ratio
        picture.crop_left = (1 - visible) / 2
        picture.crop_right = (1 - visible) / 2
    else:
        visible = image_ratio / target_ratio
        picture.crop_top = (1 - visible) / 2
        picture.crop_bottom = (1 - visible) / 2
    border = ppt_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False)
    border.fill.background()
    return picture


def add_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def add_stage(slide, label, detail, x, y, w=2.05, fill=WHITE, accent=TEAL):
    ppt_rect(slide, x, y, w, 0.78, fill=fill, line=accent)
    ppt_text(slide, label, x + 0.14, y + 0.1, w - 0.28, 0.25, 13, accent, True)
    ppt_text(slide, detail, x + 0.14, y + 0.39, w - 0.28, 0.22, 9, MUTED)


def build_pptx(output: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PptInches(PPT_W)
    presentation.slide_height = PptInches(PPT_H)
    blank = presentation.slide_layouts[6]

    # 1. Cover
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(CHARCOAL)
    ppt_rect(slide, 0.7, 0.7, 0.14, 5.95, fill=TEAL, line=TEAL, radius=False)
    ppt_text(slide, "智学工坊", 1.15, 1.2, 5.4, 0.65, 32, WHITE, True)
    ppt_text(slide, "EduAgent Studio", 1.15, 1.92, 5.8, 0.48, 23, "8DE2D8", True, font="Aptos")
    ppt_text(slide, "基于课程证据的个性化学习与多智能体资源平台", 1.15, 2.72, 5.35, 1.0, 19, "D6E0E8", True)
    ppt_text(slide, "不是问答 Demo，而是可追溯、可审核、会持续更新的学习闭环。", 1.15, 4.08, 5.3, 0.8, 14, "9EADB9")
    orbit_x = [7.55, 10.25, 10.55, 8.05, 9.0]
    orbit_y = [1.28, 1.68, 4.48, 5.1, 3.0]
    orbit_labels = ["课程证据", "学生状态", "教师审核", "学习反馈", "智能体"]
    for index, (x, y, label) in enumerate(zip(orbit_x, orbit_y, orbit_labels)):
        fill = TEAL if index in (0, 4) else (CORAL if index == 2 else "253241")
        ppt_rect(slide, x, y, 1.75, 0.72, fill=fill, line=fill)
        ppt_text(slide, label, x, y + 0.2, 1.75, 0.25, 12, WHITE, True, align=PP_ALIGN.CENTER)
    for start, end in zip(range(5), [1, 2, 3, 4, 0]):
        ppt_line(slide, orbit_x[start] + 0.88, orbit_y[start] + 0.72, orbit_x[end] + 0.88, orbit_y[end], "536779", 1.5)
    ppt_text(slide, "第十五届中国软件杯 · A3", 1.15, 6.55, 5.5, 0.25, 10, "8192A1", True, font="Aptos")
    add_notes(slide, "开场强调：项目不是单次生成工具，而是课程证据、学生状态、智能体协作、教师审核和学习反馈组成的闭环。")

    # 2. Contest requirement as loop
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 2, "赛题要求转化为一条真实闭环", "FROM REQUIREMENTS TO PRODUCT")
    problems = [
        ("通用生成", "不了解具体课程证据"),
        ("静态画像", "无法反映学习变化"),
        ("直接发布", "未经审核放大内容风险"),
    ]
    for idx, (title, detail) in enumerate(problems):
        y = 1.7 + idx * 1.25
        ppt_rect(slide, 0.7, y, 3.25, 0.92, fill=CORAL_LIGHT, line="F2B7B1")
        ppt_text(slide, title, 0.92, y + 0.13, 1.05, 0.26, 14, CORAL, True)
        ppt_text(slide, detail, 1.92, y + 0.14, 1.76, 0.48, 11, INK)
    stages = ["课程知识库", "12 维画像", "五智能体生成", "教师审核", "学习与反馈"]
    for idx, stage in enumerate(stages):
        y = 1.55 + idx * 0.94
        ppt_rect(slide, 5.0, y, 3.0, 0.6, fill=WHITE, line=TEAL)
        ppt_text(slide, stage, 5.15, y + 0.16, 2.7, 0.24, 13, TEAL, True, align=PP_ALIGN.CENTER)
        if idx < len(stages) - 1:
            ppt_line(slide, 6.5, y + 0.6, 6.5, y + 0.92, TEAL, 2)
    ppt_rect(slide, 8.65, 1.58, 3.75, 4.5, fill=CHARCOAL, line=CHARCOAL)
    ppt_text(slide, "CURRENT EVIDENCE", 8.95, 1.92, 3.1, 0.28, 10, "8DE2D8", True, font="Aptos")
    facts = [("1", "完整课程"), ("9", "章节证据"), ("12", "画像维度"), ("9", "资源类型")]
    for idx, (number, label) in enumerate(facts):
        x = 8.95 + (idx % 2) * 1.6
        y = 2.55 + (idx // 2) * 1.45
        ppt_text(slide, number, x, y, 1.2, 0.55, 27, WHITE, True, font="Aptos")
        ppt_text(slide, label, x, y + 0.6, 1.35, 0.3, 11, "AFC0CC")
    ppt_text(slide, "必做项不是页面清单，而是可持续运行的数据链。", 8.95, 5.32, 3.0, 0.5, 12, "D6E0E8", True)
    add_notes(slide, "解释三类行业问题，并用右侧闭环说明我们如何逐项落地 A3 必做能力。")

    # 3. Product panorama
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 3, "对话入口与专业工作台并存", "PRODUCT PANORAMA")
    lanes = [
        ("学生", ["Tutor", "画像", "路径", "任务", "资源", "反馈/报告"], BLUE, BLUE_LIGHT),
        ("教师", ["课程", "知识库", "智能体", "审核", "学情分析"], TEAL, TEAL_LIGHT),
        ("管理员", ["用户/课程", "模型/Prompt", "成本/安全", "审计日志"], CORAL, CORAL_LIGHT),
    ]
    for lane_idx, (role, steps, accent, light) in enumerate(lanes):
        y = 1.65 + lane_idx * 1.45
        ppt_rect(slide, 0.7, y, 1.28, 0.88, fill=accent, line=accent)
        ppt_text(slide, role, 0.7, y + 0.27, 1.28, 0.28, 15, WHITE, True, align=PP_ALIGN.CENTER)
        step_width = 9.8 / len(steps)
        for idx, step in enumerate(steps):
            x = 2.25 + idx * step_width
            ppt_rect(slide, x, y, step_width - 0.12, 0.88, fill=light, line=accent)
            ppt_text(slide, step, x + 0.05, y + 0.28, step_width - 0.22, 0.27, 11.5, accent, True, align=PP_ALIGN.CENTER)
    ppt_rect(slide, 0.7, 6.12, 11.95, 0.58, fill=CHARCOAL, line=CHARCOAL)
    ppt_text(slide, "统一数据层  ·  课程  ·  知识点  ·  画像  ·  掌握度  ·  资源  ·  证据  ·  审核  ·  反馈", 0.95, 6.28, 11.45, 0.24, 12, WHITE, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "说明产品不是把功能塞进一个首页。对话是学生主入口，教师与管理员保留高密度专业工作台。")

    # 4. Architecture
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 4, "真实实施架构：统一业务、安全与证据边界", "IMPLEMENTED ARCHITECTURE")
    layers = [
        ("React 18 SPA", "三角色路由 · Zustand · 可访问响应式界面", BLUE, BLUE_LIGHT),
        ("REST + SSE / FastAPI", "认证 · RBAC · 课程数据域 · Pydantic 合约", TEAL, TEAL_LIGHT),
        ("Service / Repository / Transaction", "LangGraph · Tutor Supervisor · LLM Gateway · Safety", "374151", "E8EDF1"),
        ("Course RAG + Evidence", "MySQL chunk · BM25 · 中文 n-gram · 引用校验", CORAL, CORAL_LIGHT),
        ("Runtime", "MySQL 8 · Redis 7 · SQLite Checkpoint · Persistent Files", CHARCOAL, "DCE3E8"),
    ]
    for idx, (title, detail, accent, light) in enumerate(layers):
        y = 1.55 + idx * 0.96
        width = 10.7 - idx * 0.42
        x = 1.3 + idx * 0.21
        ppt_rect(slide, x, y, width, 0.72, fill=light, line=accent)
        ppt_text(slide, title, x + 0.2, y + 0.13, 3.2, 0.26, 14, accent, True)
        ppt_text(slide, detail, x + 3.45, y + 0.14, width - 3.65, 0.26, 11, INK)
    ppt_text(slide, "当前链路不依赖 PostgreSQL / pgvector / MinIO", 3.55, 6.43, 6.4, 0.3, 12, CORAL, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "架构陈述严格对应当前仓库。特别说明检索基于 MySQL chunk 与 BM25，不写不存在的 pgvector。")

    # 5. Knowledge evidence
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 5, "课程知识库：生成结果可以回到原文核对", "COURSE EVIDENCE")
    flow = ["上传校验", "版本分块", "教师关联", "BM25 检索", "引用校验"]
    for idx, label in enumerate(flow):
        x = 0.68 + idx * 1.0
        ppt_rect(slide, x, 1.58, 0.86, 0.68, fill=TEAL_LIGHT, line=TEAL)
        ppt_text(slide, label, x + 0.03, 1.78, 0.8, 0.28, 10.5, TEAL, True, align=PP_ALIGN.CENTER)
        if idx < len(flow) - 1:
            ppt_line(slide, x + 0.86, 1.92, x + 1.0, 1.92, TEAL, 2)
    ppt_text(slide, "9/9", 0.75, 2.7, 2.0, 0.72, 32, CHARCOAL, True, font="Aptos")
    ppt_text(slide, "CS301 主题可召回", 0.76, 3.42, 2.7, 0.35, 12, MUTED, True)
    add_bullets(slide, ["9 个章节", "9 个知识点", "9 条 confirmed 关联", "无证据则降级为草稿"], 0.76, 4.0, 3.7, 1.85, 12)
    add_picture_crop(slide, SCREENSHOTS / "teacher-knowledge-base.png", 5.05, 2.38, 7.6, 4.41)
    add_notes(slide, "展示事务与 ACID 的真实检索结果，强调 chunk ID、原文与教师确认关联。")

    # 6. Multi-agent workflow
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 6, "五智能体资源生产：有状态、可返工、可恢复", "MULTI-AGENT PRODUCTION")
    labels = [("诊断", "薄弱点"), ("规划", "路径"), ("生成", "课程证据"), ("评测", "质量分"), ("审核建议", "教师决策")]
    for idx, (label, detail) in enumerate(labels):
        add_stage(slide, label, detail, 0.68 + idx * 1.55, 1.58, w=1.35, fill=WHITE, accent=TEAL if idx != 4 else CORAL)
        if idx < len(labels) - 1:
            ppt_line(slide, 2.03 + idx * 1.55, 1.97, 2.2 + idx * 1.55, 1.97, TEAL, 2)
    ppt_text(slide, "score < 7  →  revision（最多 3 次）", 1.7, 2.63, 5.7, 0.28, 11, CORAL, True, align=PP_ALIGN.CENTER)
    add_picture_crop(slide, SCREENSHOTS / "teacher-agent-workbench.png", 0.68, 3.08, 11.97, 3.72)
    add_notes(slide, "截图中指认节点轨迹、质量分、修订次数和生成结果。Teacher Review Agent 只给建议，最终发布由教师决定。")

    # 7. Profile and path
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 7, "12 维画像驱动动态学习路径", "PROFILE TO PATH")
    ppt_text(slide, "自然语言候选 → 学生确认 → 原子更新 → 前后快照", 0.72, 1.48, 5.75, 0.36, 13, TEAL, True)
    ppt_text(slide, "知识层级 + 前置关系 + 掌握度 → 当前学习点与资源建议", 6.78, 1.48, 5.85, 0.36, 13, BLUE, True)
    add_picture_crop(slide, SCREENSHOTS / "student-profile.png", 0.68, 1.95, 5.9, 4.82)
    add_picture_crop(slide, SCREENSHOTS / "student-learning-path.png", 6.75, 1.95, 5.9, 4.82)
    add_notes(slide, "说明个性化不是 Prompt 里一句话，而是学生确认的结构化画像和持续更新的知识点掌握度。")

    # 8. Tutor
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 8, "Tutor：问题驱动的课程证据与工具协同", "CONVERSATIONAL ENTRY")
    ppt_rect(slide, 0.68, 1.56, 3.45, 5.22, fill=CHARCOAL, line=CHARCOAL)
    ppt_text(slide, "示例问题", 0.98, 1.9, 2.8, 0.28, 11, "8DE2D8", True)
    ppt_text(slide, "用银行转账解释原子性和隔离性的区别，再给我 2 道判断题。", 0.98, 2.34, 2.72, 1.2, 18, WHITE, True)
    add_bullets(slide, ["课程检索", "证据讲解", "题型/题量约束", "SSE 会话与反馈落库"], 0.98, 4.0, 2.8, 1.8, 12, color="D6E0E8")
    ppt_text(slide, "Provider 无 tool calling → 确定性降级路由", 0.98, 6.1, 2.75, 0.42, 10.5, "8DE2D8", True)
    add_picture_crop(slide, SCREENSHOTS / "student-tutor.png", 4.38, 1.56, 8.27, 5.22)
    add_notes(slide, "现场指认用户问题、ACID 引用、两道判断题和反馈按钮。强调开发模式也执行真实课程检索。")

    # 9. Learning loop
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 9, "学习效果闭环：一次行为改变后续路径", "LEARNING OUTCOME LOOP")
    stages = [("学习行为", "任务 / 资源 / Tutor"), ("结构反馈", "测验 / 自评 / 疑问"), ("画像更新", "知识点掌握度"), ("重新推荐", "路径与资源重排")]
    for idx, (label, detail) in enumerate(stages):
        add_stage(slide, label, detail, 0.68 + idx * 1.55, 1.58, w=1.35, fill=WHITE, accent=BLUE if idx < 2 else TEAL)
        if idx < len(stages) - 1:
            ppt_line(slide, 2.03 + idx * 1.55, 1.97, 2.2 + idx * 1.55, 1.97, TEAL, 2)
    ppt_text(slide, "学生进度独立 · 只见 approved 资源 · 不能修改他人数据", 7.25, 1.72, 5.1, 0.42, 12, CORAL, True, align=PP_ALIGN.RIGHT)
    add_picture_crop(slide, SCREENSHOTS / "student-feedback.png", 0.68, 2.55, 11.97, 4.25)
    add_notes(slide, "展示最近反馈、关联资源、掌握度和下一步推荐，说明反馈如何进入画像和推荐。")

    # 10. Governance
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 10, "管理员治理：模型能力可配置、可限额、可追责", "GOVERNANCE & AUDIT")
    governance = [
        ("模型与凭证", "Provider / Model / Prompt\nAES-GCM 加密，不回显明文", TEAL, TEAL_LIGHT),
        ("成本与性能", "Token / 耗时 / 成本\n预算阈值与失败原因", BLUE, BLUE_LIGHT),
        ("内容安全", "输入 / 输出 / 工具参数\n注入与凭证暴露双检", CORAL, CORAL_LIGHT),
        ("操作审计", "登录 / 审核 / 策略变更\n角色、对象、结果可追踪", "5B6472", "E8EDF1"),
    ]
    for idx, (title, detail, accent, light) in enumerate(governance):
        x = 0.7 + (idx % 2) * 6.05
        y = 1.62 + (idx // 2) * 2.2
        ppt_rect(slide, x, y, 5.75, 1.72, fill=light, line=accent)
        ppt_text(slide, title, x + 0.25, y + 0.22, 2.2, 0.3, 16, accent, True)
        ppt_text(slide, detail, x + 0.25, y + 0.72, 5.1, 0.62, 12, INK)
    ppt_rect(slide, 0.7, 6.14, 11.95, 0.55, fill=CHARCOAL, line=CHARCOAL)
    ppt_text(slide, "生产环境强制真实 Provider · 强密钥 · 非通配 CORS · 禁止 Mock", 0.95, 6.29, 11.45, 0.25, 12, WHITE, True, align=PP_ALIGN.CENTER)
    add_notes(slide, "当前页使用可编辑治理矩阵，避免在没有管理员截图的情况下伪造界面。答辩时可切到真实管理端展示。")

    # 11. Verification
    slide = presentation.slides.add_slide(blank)
    add_slide_chrome(slide, 11, "工程验证结果来自测试、构建和真实运行", "ENGINEERING EVIDENCE")
    metrics = [("271", "后端测试", TEAL), ("9/9", "课程主题召回", BLUE), ("32", "三角色主路由", CORAL), ("0", "npm 生产漏洞", "475569")]
    for idx, (number, label, accent) in enumerate(metrics):
        x = 0.7 + idx * 3.0
        ppt_rect(slide, x, 1.68, 2.72, 1.75, fill=WHITE, line=accent)
        ppt_text(slide, number, x, 1.92, 2.72, 0.68, 31, accent, True, font="Aptos", align=PP_ALIGN.CENTER)
        ppt_text(slide, label, x, 2.72, 2.72, 0.28, 11.5, MUTED, True, align=PP_ALIGN.CENTER)
    checks = [
        ("响应式", "390 × 844 / 1024 / 1440"),
        ("前端门禁", "TypeScript + Vite build"),
        ("依赖一致性", "npm audit 0 / pip check"),
        ("运行健康", "API + MySQL + Redis + Web"),
    ]
    for idx, (label, detail) in enumerate(checks):
        y = 4.05 + idx * 0.57
        ppt_text(slide, label, 1.0, y, 2.0, 0.26, 12, CHARCOAL, True)
        ppt_rect(slide, 3.15, y - 0.03, 8.8, 0.38, fill="E7EDF0", line="E7EDF0")
        ppt_text(slide, detail, 3.35, y + 0.04, 8.45, 0.2, 10.5, MUTED, True)
    add_notes(slide, "明确 271 项测试是当前全量结果。外部讯飞凭证联调和最终 MP4 不计入已完成基线。")

    # 12. Closing
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(CHARCOAL)
    ppt_text(slide, "让 AI 教学从“生成一次”\n变成“持续闭环”", 0.85, 0.8, 5.2, 1.45, 30, WHITE, True)
    differences = [
        "学生确认的 12 维画像",
        "真实 chunk 与知识点证据",
        "有状态多智能体与质量返工",
        "对话入口 + 专业工作台",
        "教师保留最终发布决策",
        "成本、安全与操作统一审计",
    ]
    for idx, item in enumerate(differences):
        x = 0.9 + (idx % 2) * 5.95
        y = 2.65 + (idx // 2) * 0.92
        accent = TEAL if idx % 2 == 0 else CORAL
        ppt_rect(slide, x, y, 5.55, 0.68, fill="253241", line=accent)
        ppt_text(slide, f"{idx + 1:02d}", x + 0.18, y + 0.2, 0.5, 0.23, 11, accent, True, font="Aptos")
        ppt_text(slide, item, x + 0.78, y + 0.18, 4.45, 0.28, 13, WHITE, True)
    ppt_text(slide, "交付边界", 0.9, 5.83, 1.4, 0.28, 11, "8DE2D8", True)
    ppt_text(slide, "讯飞真实能力需有效凭证联调；最终 MP4 按当前脚本录制；pip-audit 需在具备工具的环境执行。", 2.25, 5.78, 10.0, 0.55, 11.5, "CAD6DE")
    ppt_text(slide, "EduAgent Studio", 0.9, 6.78, 4.2, 0.3, 11, "8192A1", True, font="Aptos")
    add_notes(slide, "收束差异化，并主动说明真实外部 API、MP4 和 pip-audit 的交付边界。")

    presentation.core_properties.title = "智学工坊 EduAgent Studio 参赛演示"
    presentation.core_properties.subject = "中国软件杯 A3 赛题"
    presentation.core_properties.author = "EduAgent Studio 参赛团队"
    presentation.core_properties.keywords = "EduAgent Studio, A3, 多智能体, 个性化学习"
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--docs-only", action="store_true")
    parser.add_argument("--ppt-only", action="store_true")
    args = parser.parse_args()

    if not args.ppt_only:
        sources = [
            (SUBMISSION / "docs" / "AI_Coding工具使用说明.md", SUBMISSION / "AI_Coding工具使用说明.docx", "AI Coding 工具使用说明"),
            (SUBMISSION / "docs" / "测试说明书.md", SUBMISSION / "测试说明书.docx", "软件测试说明书"),
            (SUBMISSION / "docs" / "系统开发说明书.md", SUBMISSION / "系统开发说明书.docx", "系统开发说明书"),
        ]
        for source, output, category in sources:
            build_docx(source, output, category)
            print(f"built {output.relative_to(ROOT)}")

    if not args.docs_only:
        output = SUBMISSION / "软件杯参赛演示PPT.pptx"
        build_pptx(output)
        print(f"built {output.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
